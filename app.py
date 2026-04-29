import streamlit as st
import os
import numpy as np
import pandas as pd
import nltk
import re
import requests
import json
import base64
import hashlib
import random
from io import BytesIO
from datetime import datetime
from collections import Counter
from dotenv import load_dotenv
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_community.embeddings import HuggingFaceEmbeddings
import plotly.express as px
from PIL import Image
from nltk.stem import WordNetLemmatizer
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import simpleSplit

# Download NLTK lemmatizer
@st.cache_resource
def download_nltk_lemmatizer():
    nltk.download('wordnet', quiet=True)
    nltk.download('omw-1.4', quiet=True)
    return WordNetLemmatizer()

lemmatizer = download_nltk_lemmatizer()

st.set_page_config(page_title="Contextual Reader", page_icon="📚", layout="wide")

# ---------- Helper functions ----------
def get_api_key():
    try:
        return st.secrets["GOOGLE_API_KEY"]
    except:
        return os.getenv("GOOGLE_API_KEY")

@st.cache_resource
def load_nltk_data():
    nltk.download('punkt', quiet=True)
    nltk.download('punkt_tab', quiet=True)
    nltk.download('stopwords', quiet=True)
    from nltk.corpus import stopwords
    return stopwords.words('english')

stop_words = load_nltk_data()

@st.cache_resource
def get_embeddings():
    return HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# ---------- Library persistence ----------
LIBRARY_FILE = "book_library.json"

def load_library():
    if os.path.exists(LIBRARY_FILE):
        with open(LIBRARY_FILE, "r") as f:
            return json.load(f)
    return []

def save_library(library):
    with open(LIBRARY_FILE, "w") as f:
        json.dump(library, f, indent=2)

def add_book_to_library(filename, word_count, stats, cover_b64=None, cover_url=None, status="unread", progress=0):
    library = load_library()
    for book in library:
        if book["filename"] == filename:
            return
    library.append({
        "filename": filename,
        "upload_date": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "word_count": word_count,
        "stats": stats,
        "cover_b64": cover_b64,
        "cover_url": cover_url,
        "status": status,
        "progress": progress
    })
    save_library(library)

def update_book_status(filename, status=None, progress=None, cover_b64=None, cover_url=None):
    library = load_library()
    found = False
    for book in library:
        if book["filename"] == filename:
            if status is not None:
                book["status"] = status
            if progress is not None:
                book["progress"] = progress
            if cover_b64 is not None:
                book["cover_b64"] = cover_b64
                book["cover_url"] = None
            if cover_url is not None:
                book["cover_url"] = cover_url
                book["cover_b64"] = None
            found = True
            break
    if found:
        save_library(library)
    return found

def delete_book_from_library(filename):
    library = load_library()
    library = [b for b in library if b["filename"] != filename]
    save_library(library)

# ---------- Saved summaries & quizzes persistence ----------
SAVED_DATA_FILE = "saved_data.json"

def load_saved_data():
    if os.path.exists(SAVED_DATA_FILE):
        with open(SAVED_DATA_FILE, "r") as f:
            data = json.load(f)
            return data.get("summaries", []), data.get("quizzes", [])
    return [], []

def save_saved_data(summaries, quizzes):
    with open(SAVED_DATA_FILE, "w") as f:
        json.dump({"summaries": summaries, "quizzes": quizzes}, f, indent=2)

# ---------- Definition lookup ----------
def get_definition(word):
    base_word = lemmatizer.lemmatize(word.lower())
    for attempt in [word.lower(), base_word]:
        try:
            url = f"https://api.dictionaryapi.dev/api/v2/entries/en/{attempt}"
            resp = requests.get(url, timeout=5)
            if resp.status_code == 200:
                data = resp.json()
                definition = data[0]['meanings'][0]['definitions'][0]['definition']
                return definition
        except:
            continue
    return None

# ---------- Study aids (PPT + PDF) ----------
def create_ppt(summary_text, filename="study_aids.pptx"):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Chapter Summary"
    content = slide.placeholders[1]
    lines = summary_text.split('\n')
    content.text = "\n".join(lines[:10])
    prs.save(filename)
    return filename

def create_pdf(summary_text, filename="summary.pdf"):
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    y = height - 50
    lines = simpleSplit(summary_text, 'Helvetica', 12, width - 100)
    for line in lines:
        if y < 50:
            c.showPage()
            y = height - 50
        c.drawString(50, y, line)
        y -= 15
    c.save()
    buffer.seek(0)
    return buffer

# ---------- Vector store ----------
class SimpleVectorStore:
    def __init__(self, embeddings_model):
        self.embeddings_model = embeddings_model
        self.chunks = []
        self.embeddings = np.array([])

    def add_documents(self, chunks):
        self.chunks = chunks
        texts = [c.page_content for c in chunks]
        self.embeddings = np.array(self.embeddings_model.embed_documents(texts))

    def similarity_search(self, query, k=12):
        q_emb = np.array(self.embeddings_model.embed_query(query))
        sim = np.dot(self.embeddings, q_emb) / (np.linalg.norm(self.embeddings, axis=1) * np.linalg.norm(q_emb))
        idx = np.argsort(sim)[-k:][::-1]
        return [self.chunks[i] for i in idx]

# ---------- Session state (multi‑book) ----------
if "books" not in st.session_state:
    st.session_state.books = {}   # key = filename, value = dict with vectorstore, full_text, stats, chunks_metadata, rare_words
if "active_book" not in st.session_state:
    st.session_state.active_book = None
if "saved_qa" not in st.session_state:
    st.session_state.saved_qa = []
if "saved_summaries" not in st.session_state:
    st.session_state.saved_summaries = []
if "saved_quizzes" not in st.session_state:
    st.session_state.saved_quizzes = []
if "last_question" not in st.session_state:
    st.session_state.last_question = ""
if "last_answer" not in st.session_state:
    st.session_state.last_answer = ""
if "last_context_docs" not in st.session_state:
    st.session_state.last_context_docs = []
if "custom_word_lookup" not in st.session_state:
    st.session_state.custom_word_lookup = ""
if "quiz_questions" not in st.session_state:
    st.session_state.quiz_questions = []
if "summary_cache" not in st.session_state:
    st.session_state.summary_cache = {}
if "cover_url_temp" not in st.session_state:
    st.session_state.cover_url_temp = ""
if "imported" not in st.session_state:
    st.session_state.imported = False
if "current_summary" not in st.session_state:
    st.session_state.current_summary = None
if "current_summary_source" not in st.session_state:
    st.session_state.current_summary_source = None
if "current_summary_original" not in st.session_state:
    st.session_state.current_summary_original = None

# Load persisted summaries and quizzes
saved_summaries, saved_quizzes = load_saved_data()
if saved_summaries:
    st.session_state.saved_summaries = saved_summaries
if saved_quizzes:
    st.session_state.saved_quizzes = saved_quizzes

# ---------- Sidebar (always visible) ----------
with st.sidebar:
    st.markdown("## 📖 Contextual Reader")
    st.markdown("---")

    # Book selection dropdown (if multiple books)
    if len(st.session_state.books) > 0:
        book_options = list(st.session_state.books.keys())
        if st.session_state.active_book not in book_options:
            st.session_state.active_book = book_options[0] if book_options else None
        selected_book = st.selectbox("Select active book", book_options, index=book_options.index(st.session_state.active_book) if st.session_state.active_book in book_options else 0)
        if selected_book != st.session_state.active_book:
            st.session_state.active_book = selected_book
            st.rerun()
    
    # Import library
    uploaded_lib = st.file_uploader("📥 Import library (JSON)", type="json", key="lib_uploader")
    if uploaded_lib and not st.session_state.imported:
        try:
            new_lib = json.load(uploaded_lib)
            if isinstance(new_lib, list) and all(isinstance(b, dict) and "filename" in b for b in new_lib):
                save_library(new_lib)
                st.session_state.imported = True
                st.success("Library imported! Refreshing...")
                st.rerun()
            else:
                st.error("Invalid JSON structure.")
        except Exception as e:
            st.error(f"Invalid JSON file: {e}")
    if not uploaded_lib and st.session_state.imported:
        st.session_state.imported = False
    
    if st.button("📤 Export library (JSON)"):
        lib = load_library()
        if lib:
            json_str = json.dumps(lib, indent=2)
            st.download_button("Download library", json_str, "book_library.json", "application/json")
        else:
            st.info("No library to export.")
    st.caption("ℹ️ Library JSON stores only metadata (titles, progress, covers). To keep your RAG state, you must re-upload the PDF after import.")
    
    st.markdown("---")
    if st.button("🗑️ Clear current book"):
        if st.session_state.active_book and st.session_state.active_book in st.session_state.books:
            del st.session_state.books[st.session_state.active_book]
            if st.session_state.books:
                st.session_state.active_book = list(st.session_state.books.keys())[0]
            else:
                st.session_state.active_book = None
            st.rerun()
    
    if st.button("🗑️ Clear entire library"):
        save_library([])
        st.success("All books removed from library.")
        st.rerun()

# ---------- Main area ----------
st.title("📚 Contextual Reader")
st.markdown("*Your personal AI reading companion*")

uploaded_file = st.file_uploader("Upload a PDF document", type="pdf")

if uploaded_file is not None:
    filename = uploaded_file.name
    if filename not in st.session_state.books:
        with st.spinner(f"Processing {filename}..."):
            with open("temp.pdf", "wb") as f:
                f.write(uploaded_file.getbuffer())
            loader = PyPDFLoader("temp.pdf")
            documents = loader.load()
            full_text = " ".join([d.page_content for d in documents])

            # Chunks with page numbers
            chunks_meta = []
            for doc in documents:
                splits = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100).split_text(doc.page_content)
                for s in splits:
                    chunks_meta.append((s, doc.metadata.get("page", 0)+1))

            text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
            chunks = text_splitter.split_documents(documents)
            embeddings = get_embeddings()
            vectorstore = SimpleVectorStore(embeddings)
            vectorstore.add_documents(chunks)

            # Word statistics
            words = re.findall(r'\b[a-zA-Z]{2,}\b', full_text.lower())
            unique_words = set(words)
            sentences = nltk.sent_tokenize(full_text)
            avg_sentence_len = np.mean([len(w.split()) for w in sentences]) if sentences else 0
            reading_time_min = len(words) / 250
            stats = {
                "characters": len(full_text),
                "words": len(words),
                "unique_words": len(unique_words),
                "sentences": len(sentences),
                "avg_sentence_len": round(avg_sentence_len, 1),
                "reading_time": round(reading_time_min, 1),
            }

            # Rare words
            freq = Counter(words)
            rare = {w: c for w, c in freq.items() if w not in stop_words and 1 <= c <= 3 and len(w) >= 3}
            rare_words = sorted(rare.items(), key=lambda x: x[1])[:50]

            # Store in books dict
            st.session_state.books[filename] = {
                "vectorstore": vectorstore,
                "full_text": full_text,
                "chunks_metadata": chunks_meta,
                "stats": stats,
                "rare_words": rare_words
            }
            if st.session_state.active_book is None:
                st.session_state.active_book = filename

            # Add to library JSON (only metadata)
            add_book_to_library(filename, len(words), stats)
            os.remove("temp.pdf")
            st.success(f"✅ Processed: {filename}")
            st.balloons()
    else:
        st.info(f"📚 '{filename}' is already loaded. Select it from the sidebar to ask questions.")

# ---------- Tabs (only when a book is active) ----------
if st.session_state.active_book is not None and st.session_state.active_book in st.session_state.books:
    book_data = st.session_state.books[st.session_state.active_book]
    vectorstore = book_data["vectorstore"]
    full_text = book_data["full_text"]
    chunks_metadata = book_data["chunks_metadata"]
    book_stats = book_data["stats"]
    rare_words = book_data["rare_words"]

    api_key = get_api_key()
    if not api_key:
        st.error("❌ GOOGLE_API_KEY missing.")
        st.stop()
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.7, google_api_key=api_key)

    tabs = st.tabs(["💬 Ask", "📊 Dashboard", "📌 Saved Q&A", "📖 Vocabulary", "🧠 Quiz", "📚 Library", "✍️ Summarize", "📑 Study Aids", "💾 Saved Summaries & Quizzes"])

    # ---------- TAB 1: Ask ----------
    with tabs[0]:
        st.header(f"Ask about '{st.session_state.active_book}' (literary critique)")
        question = st.text_input("Your question:", placeholder="e.g., What is your opinion of Mr. Darcy?", key="ask_input")
        ask_button = st.button("Ask", key="ask_submit")
        
        if ask_button and question:
            with st.spinner("Thinking..."):
                try:
                    docs = vectorstore.similarity_search(question, k=12)
                    if not docs:
                        st.warning("No relevant passages found.")
                        st.session_state.last_answer = ""
                        st.session_state.last_context_docs = []
                    else:
                        context = "\n\n".join([d.page_content for d in docs])
                        prompt = f"""You are a literary critic. Provide your opinion based only on the context below. Feel free to analyse characters, themes, and style. If the answer is not in the context, say "I don't have enough information from the book."

Context:
{context}

Question: {question}

Answer:"""
                        response = llm.invoke(prompt)
                        st.session_state.last_answer = response.content
                        st.session_state.last_context_docs = docs
                        st.session_state.last_question = question
                    st.rerun()
                except Exception as e:
                    st.error(f"Error: {e}")
        
        if st.session_state.last_answer:
            st.success(st.session_state.last_answer)
            if st.button("💾 Save this Q&A"):
                st.session_state.saved_qa.append({
                    "question": st.session_state.last_question,
                    "answer": st.session_state.last_answer,
                    "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                })
                st.toast("Saved!", icon="✅")
            with st.expander("📖 Source passages"):
                for i, d in enumerate(st.session_state.last_context_docs):
                    st.caption(f"Source {i+1}: {d.page_content[:500]}...")

    # ---------- TAB 2: Dashboard ----------
    with tabs[1]:
        st.header(f"Dashboard: {st.session_state.active_book}")
        if book_stats:
            s = book_stats
            cols = st.columns(4)
            cols[0].metric("Characters", f"{s['characters']:,}")
            cols[1].metric("Words", f"{s['words']:,}")
            cols[2].metric("Unique words", f"{s['unique_words']:,}")
            cols[3].metric("Reading time (min)", s['reading_time'])
            st.subheader("Character mentions")
            default = ["Elizabeth", "Darcy", "Jane", "Bingley", "Wickham", "Collins", "Lydia"]
            chars = st.text_input("Character names (comma separated)", value=",".join(default))
            if chars:
                names = [n.strip() for n in chars.split(",")]
                counts = {}
                text_lower = full_text.lower()
                for n in names:
                    pattern = r'\b' + re.escape(n.lower()) + r'\b'
                    counts[n] = len(re.findall(pattern, text_lower))
                fig = px.bar(x=list(counts.keys()), y=list(counts.values()), color_discrete_sequence=["#4a6fa5"])
                st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("No stats available.")

    # ---------- TAB 3: Saved Q&A ----------
    with tabs[2]:
        st.header("Saved Q&A")
        if st.session_state.saved_qa:
            df = pd.DataFrame(st.session_state.saved_qa)
            st.dataframe(df)
            csv = df.to_csv(index=False).encode()
            st.download_button("Download CSV", csv, "saved_qa.csv")
        else:
            st.info("No saved Q&A yet.")

    # ---------- TAB 4: Vocabulary ----------
    with tabs[3]:
        st.header("Vocabulary Builder")
        if rare_words:
            rare_df = pd.DataFrame(rare_words, columns=["Word", "Frequency"])
            st.dataframe(rare_df, use_container_width=True)
            selected_word = st.selectbox("Select a rare word to see its definition", [w for w, _ in rare_words])
            if selected_word:
                if selected_word not in st.session_state.word_definitions:
                    with st.spinner("Looking up definition..."):
                        defn = get_definition(selected_word)
                        st.session_state.word_definitions[selected_word] = defn if defn else "Definition not found."
                st.success(f"**{selected_word}**: {st.session_state.word_definitions[selected_word]}")
        else:
            st.info("No rare words found in this book.")

        st.markdown("---")
        st.subheader("Look up any word")
        custom_word = st.text_input("Enter any word", value=st.session_state.custom_word_lookup, key="custom_word_input")
        if st.button("Get definition"):
            if custom_word:
                st.session_state.custom_word_lookup = custom_word
                if custom_word not in st.session_state.word_definitions:
                    with st.spinner("Looking up definition..."):
                        defn = get_definition(custom_word)
                        st.session_state.word_definitions[custom_word] = defn if defn else "Definition not found."
                st.success(f"**{custom_word}**: {st.session_state.word_definitions[custom_word]}")
            else:
                st.warning("Please enter a word.")

        if st.button("📤 Export rare words + definitions (CSV)"):
            export_data = []
            for word, freq in rare_words:
                if word not in st.session_state.word_definitions:
                    defn = get_definition(word)
                    st.session_state.word_definitions[word] = defn if defn else "Definition not found."
                export_data.append({"Word": word, "Frequency": freq, "Definition": st.session_state.word_definitions.get(word, "Not available")})
            export_df = pd.DataFrame(export_data)
            csv = export_df.to_csv(index=False).encode()
            st.download_button("Download CSV", csv, "vocabulary.csv", "text/csv")

    # ---------- TAB 5: Quiz ----------
    with tabs[4]:
        st.header("AI‑Generated Quiz")
        col1, col2 = st.columns([2,1])
        with col1:
            if st.button("🎲 Generate new quiz (5 questions)"):
                with st.spinner("Generating quiz..."):
                    try:
                        total_chunks = len(vectorstore.chunks)
                        if total_chunks == 0:
                            st.error("No text chunks available.")
                        else:
                            sample_size = min(20, total_chunks)
                            random_indices = random.sample(range(total_chunks), sample_size)
                            sample_chunks = [vectorstore.chunks[i].page_content for i in random_indices]
                            context = "\n\n".join(sample_chunks)[:8000]
                            prompt = f"""Generate 5 multiple-choice questions based on the following text. Return a JSON list with each question containing: "question", "options" (list of 4 strings), "correct" (0-based index of correct option). Example: [{{"question": "What is X?", "options": ["A", "B", "C", "D"], "correct": 1}}]

Text:
{context}"""
                            response = llm.invoke(prompt)
                            text = response.content
                            start = text.find('[')
                            end = text.rfind(']')+1
                            if start != -1 and end != -1:
                                json_str = text[start:end]
                                st.session_state.quiz_questions = json.loads(json_str)
                                st.success("Quiz generated!")
                            else:
                                st.error("Could not parse quiz. Try again.")
                    except Exception as e:
                        st.error(f"Error: {e}")
        with col2:
            if st.session_state.quiz_questions:
                if st.button("💾 Save this quiz"):
                    st.session_state.saved_quizzes.append({
                        "questions": st.session_state.quiz_questions,
                        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "book": st.session_state.active_book
                    })
                    save_saved_data(st.session_state.saved_summaries, st.session_state.saved_quizzes)
                    st.toast("Quiz saved!", icon="✅")

        if st.session_state.quiz_questions:
            with st.form(key="quiz_form"):
                user_answers = []
                for i, q in enumerate(st.session_state.quiz_questions):
                    st.markdown(f"**Q{i+1}: {q['question']}**")
                    selected = st.radio(f"Select answer", q['options'], key=f"quiz_radio_{i}", index=None)
                    user_answers.append(selected)
                submitted = st.form_submit_button("Submit answers")
                if submitted:
                    correct = 0
                    for i, (q, ans) in enumerate(zip(st.session_state.quiz_questions, user_answers)):
                        if ans == q['options'][q['correct']]:
                            correct += 1
                            st.success(f"Q{i+1}: Correct!")
                        else:
                            st.error(f"Q{i+1}: Incorrect. Correct: {q['options'][q['correct']]}")
                    st.info(f"Score: {correct}/{len(st.session_state.quiz_questions)}")
        else:
            st.info("Click 'Generate new quiz' to start.")

    # ---------- TAB 6: Library ----------
    with tabs[5]:
        st.header("📚 Your Library")
        library = load_library()
        if not library:
            st.info("No books in library. Upload a PDF to add it.")
        else:
            total_books = len(library)
            total_words = sum(b.get("word_count", 0) for b in library)
            finished = sum(1 for b in library if b.get("status") == "finished")
            reading = sum(1 for b in library if b.get("status") == "reading")
            rate = (finished / total_books * 100) if total_books else 0
            
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("📚 Total books", total_books)
            c2.metric("📝 Total words", f"{total_words:,}")
            c3.metric("✅ Finished", finished)
            c4.metric("📖 Reading now", reading)
            st.progress(rate/100, text=f"Completion: {rate:.1f}%")
            st.markdown("---")
            
            sort_by = st.selectbox("Sort by", ["Upload date (newest first)", "Word count", "Title", "Status", "Progress"])
            if sort_by == "Upload date (newest first)":
                library.sort(key=lambda x: x.get("upload_date", ""), reverse=True)
            elif sort_by == "Word count":
                library.sort(key=lambda x: x.get("word_count", 0), reverse=True)
            elif sort_by == "Title":
                library.sort(key=lambda x: x.get("filename", "").lower())
            elif sort_by == "Status":
                library.sort(key=lambda x: x.get("status", ""))
            elif sort_by == "Progress":
                library.sort(key=lambda x: x.get("progress", 0), reverse=True)
            
            for idx, book in enumerate(library):
                with st.container():
                    cols = st.columns([1, 2, 2, 1, 1])
                    with cols[0]:
                        if book.get("cover_b64"):
                            try:
                                st.image(base64.b64decode(book["cover_b64"]), width=100)
                            except:
                                st.image("https://via.placeholder.com/100x150?text=Error", width=100)
                        elif book.get("cover_url"):
                            st.image(book["cover_url"], width=100)
                        else:
                            st.image("https://via.placeholder.com/100x150?text=No+Cover", width=100)
                    with cols[1]:
                        st.markdown(f"**{book['filename']}**")
                        st.caption(f"Uploaded: {book['upload_date']}")
                        st.caption(f"Words: {book['word_count']:,}")
                    with cols[2]:
                        new_cover_img = st.file_uploader("Update cover", type=["jpg","png"], key=f"cover_{idx}_{book['filename']}", label_visibility="collapsed")
                        new_cover_url = st.text_input("Or URL", key=f"coverurl_{idx}_{book['filename']}", placeholder="https://...")
                        if st.button("Set cover", key=f"setcover_{idx}_{book['filename']}"):
                            if new_cover_img:
                                img = Image.open(new_cover_img)
                                buffered = BytesIO()
                                img.save(buffered, format="JPEG")
                                cover_b64 = base64.b64encode(buffered.getvalue()).decode()
                                update_book_status(book['filename'], cover_b64=cover_b64)
                                st.success("Cover updated!")
                                st.rerun()
                            elif new_cover_url:
                                update_book_status(book['filename'], cover_url=new_cover_url)
                                st.success("Cover URL updated!")
                                st.rerun()
                    with cols[3]:
                        status_opts = ["unread", "reading", "finished"]
                        status_idx = status_opts.index(book.get("status", "unread"))
                        new_status = st.selectbox("Status", status_opts, index=status_idx, key=f"status_{idx}_{book['filename']}", label_visibility="collapsed")
                        new_progress = st.number_input("Pages read", min_value=0, max_value=10000, value=book.get("progress",0), key=f"progress_{idx}_{book['filename']}", label_visibility="collapsed")
                        if st.button("Update", key=f"update_{idx}_{book['filename']}"):
                            update_book_status(book['filename'], new_status, new_progress)
                            st.rerun()
                    with cols[4]:
                        if st.button("Delete", key=f"del_{idx}_{book['filename']}"):
                            delete_book_from_library(book['filename'])
                            st.rerun()
                    st.divider()

    # ---------- TAB 7: Summarize ----------
    with tabs[6]:
        st.header("Summarize Chapters, Paragraphs, or Custom Text")
        st.markdown("**Note:** Summarization calls the Gemini API each time (spinner appears). You can save summaries for later.")
        
        if chunks_metadata:
            pages = sorted(set([page for _, page in chunks_metadata]))
            min_page = min(pages) if pages else 1
            max_page = max(pages) if pages else 1
            page_range = st.slider("Select page range from the book", min_value=min_page, max_value=max_page, value=(min_page, min(min_page+10, max_page)), key="page_range_slider")
            
            if st.button("Summarize selected pages", key="summarize_pages"):
                selected_text = []
                for text, page in chunks_metadata:
                    if page_range[0] <= page <= page_range[1]:
                        selected_text.append(text)
                if selected_text:
                    text_to_sum = "\n\n".join(selected_text)[:8000]
                    hash_key = hashlib.md5(text_to_sum.encode()).hexdigest()
                    if hash_key in st.session_state.summary_cache:
                        summary = st.session_state.summary_cache[hash_key]
                    else:
                        with st.spinner("Generating summary..."):
                            prompt = f"Summarize the following text concisely:\n\n{text_to_sum}"
                            response = llm.invoke(prompt)
                            summary = response.content
                            st.session_state.summary_cache[hash_key] = summary
                    st.session_state.current_summary = summary
                    st.session_state.current_summary_source = f"Pages {page_range[0]}-{page_range[1]}"
                    st.session_state.current_summary_original = text_to_sum[:200] + "..."
                    st.rerun()
                else:
                    st.warning("No text in that page range.")
            
            if st.session_state.current_summary and st.session_state.current_summary_source.startswith("Pages"):
                st.success(st.session_state.current_summary)
                if st.button("💾 Save this summary", key="save_summary_pages"):
                    st.session_state.saved_summaries.append({
                        "original": st.session_state.current_summary_original,
                        "summary": st.session_state.current_summary,
                        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "source": st.session_state.current_summary_source
                    })
                    save_saved_data(st.session_state.saved_summaries, st.session_state.saved_quizzes)
                    st.toast("Summary saved!", icon="✅")
        else:
            st.info("No book loaded. Upload a PDF to summarize its pages.")
        
        st.markdown("---")
        st.subheader("Or paste your own text")
        manual_text = st.text_area("Paste text to summarize", height=200, key="custom_summary_text")
        
        if st.button("Summarize custom text", key="summarize_custom"):
            if manual_text:
                text_to_sum = manual_text[:8000]
                hash_key = hashlib.md5(text_to_sum.encode()).hexdigest()
                if hash_key in st.session_state.summary_cache:
                    summary = st.session_state.summary_cache[hash_key]
                else:
                    with st.spinner("Generating summary..."):
                        prompt = f"Summarize the following text concisely:\n\n{text_to_sum}"
                        response = llm.invoke(prompt)
                        summary = response.content
                        st.session_state.summary_cache[hash_key] = summary
                st.session_state.current_summary = summary
                st.session_state.current_summary_source = "Custom text"
                st.session_state.current_summary_original = manual_text[:200] + "..."
                st.rerun()
            else:
                st.warning("Please enter some text.")
        
        if st.session_state.current_summary and st.session_state.current_summary_source == "Custom text":
            st.success(st.session_state.current_summary)
            if st.button("💾 Save this summary", key="save_summary_custom"):
                st.session_state.saved_summaries.append({
                    "original": st.session_state.current_summary_original,
                    "summary": st.session_state.current_summary,
                    "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "source": "Custom text"
                })
                save_saved_data(st.session_state.saved_summaries, st.session_state.saved_quizzes)
                st.toast("Summary saved!", icon="✅")

    # ---------- TAB 8: Study Aids ----------
    with tabs[7]:
        st.header("Generate Study Aids")
        st.markdown("Create a PowerPoint summary and a PDF handout from a textbook chapter.")
        
        if chunks_metadata:
            pages = sorted(set([page for _, page in chunks_metadata]))
            min_page = min(pages) if pages else 1
            max_page = max(pages) if pages else 1
            study_range = st.slider("Select page range for study aids", min_value=min_page, max_value=max_page, value=(min_page, min(min_page+10, max_page)), key="study_range")
            
            if st.button("Generate Study Aids", key="gen_study"):
                selected_text = []
                for text, page in chunks_metadata:
                    if study_range[0] <= page <= study_range[1]:
                        selected_text.append(text)
                if selected_text:
                    study_text = "\n\n".join(selected_text)[:8000]
                    with st.spinner("Generating summary and slides..."):
                        prompt = f"Summarize the following text in 5 bullet points suitable for a PowerPoint slide:\n\n{study_text}"
                        response = llm.invoke(prompt)
                        summary_text = response.content
                        st.subheader("Chapter Summary")
                        st.write(summary_text)
                        
                        # Create PPT
                        try:
                            ppt_filename = create_ppt(summary_text, "study_aids.pptx")
                            with open(ppt_filename, "rb") as f:
                                st.download_button("Download PowerPoint (.pptx)", f, "study_aids.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                            os.remove(ppt_filename)
                        except Exception as e:
                            st.error(f"PPT generation error: {e}")
                        
                        # Create PDF
                        try:
                            pdf_buffer = create_pdf(summary_text, "summary.pdf")
                            st.download_button("Download PDF Summary", pdf_buffer, "chapter_summary.pdf", "application/pdf")
                        except Exception as e:
                            st.error(f"PDF generation error: {e}")
                        
                        if st.button("💾 Save this summary to my saved summaries", key="save_study_summary"):
                            st.session_state.saved_summaries.append({
                                "original": study_text[:200]+"...",
                                "summary": summary_text,
                                "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                                "source": "Study aids generation"
                            })
                            save_saved_data(st.session_state.saved_summaries, st.session_state.saved_quizzes)
                            st.toast("Saved!")
                else:
                    st.warning("No text in that page range.")
        else:
            st.info("No book loaded. Upload a PDF to generate study aids.")

    # ---------- TAB 9: Saved Summaries & Quizzes ----------
    with tabs[8]:
        st.header("Saved Summaries & Quizzes")
        
        st.subheader("📝 Saved Summaries")
        if st.session_state.saved_summaries:
            for i, s in enumerate(st.session_state.saved_summaries):
                with st.expander(f"Summary {i+1} - {s['timestamp']} (Source: {s['source']})"):
                    st.caption(f"Original excerpt: {s['original']}")
                    st.write(f"**Summary:** {s['summary']}")
            if st.button("Export all summaries as JSON"):
                json_str = json.dumps(st.session_state.saved_summaries, indent=2)
                st.download_button("Download", json_str, "saved_summaries.json", "application/json")
        else:
            st.info("No saved summaries yet.")
        
        st.markdown("---")
        st.subheader("📋 Saved Quizzes")
        if st.session_state.saved_quizzes:
            for i, qz in enumerate(st.session_state.saved_quizzes):
                with st.expander(f"Quiz {i+1} - {qz['timestamp']} (Book: {qz['book']})"):
                    for j, q in enumerate(qz['questions']):
                        st.write(f"**Q{j+1}:** {q['question']}")
                        for opt in q['options']:
                            st.write(f"  - {opt}")
                        st.write(f"*Correct answer:* {q['options'][q['correct']]}")
            if st.button("Export all quizzes as JSON"):
                json_str = json.dumps(st.session_state.saved_quizzes, indent=2)
                st.download_button("Download", json_str, "saved_quizzes.json", "application/json")
        else:
            st.info("No saved quizzes yet.")

else:
    st.info("👈 Upload a PDF to start, or select an existing book from the sidebar.")